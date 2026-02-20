"""Tests for powerpoint_tool - Generate PowerPoint files (.pptx) schema-first."""

import importlib.util
from pathlib import Path
from unittest.mock import patch

import pytest
from fastmcp import FastMCP
from pptx import Presentation

pptx_available = importlib.util.find_spec("pptx") is not None

pytestmark = pytest.mark.skipif(not pptx_available, reason="python-pptx not installed")

if pptx_available:
    from aden_tools.tools.powerpoint_tool.powerpoint_tool import register_tools

TEST_WORKSPACE_ID = "test-workspace"
TEST_AGENT_ID = "test-agent"
TEST_SESSION_ID = "test-session"


@pytest.fixture
def powerpoint_tools(mcp: FastMCP, tmp_path: Path):
    """Register PowerPoint tools and return callable tool handles."""
    with patch("aden_tools.tools.file_system_toolkits.security.WORKSPACES_DIR", str(tmp_path)):
        register_tools(mcp)
        yield {
            "powerpoint_generate": mcp._tool_manager._tools["powerpoint_generate"].fn,
            "tmp_path": tmp_path,
        }


def test_powerpoint_generate_creates_pptx(powerpoint_tools):
    gen = powerpoint_tools["powerpoint_generate"]

    deck = {
        "title": "Weekly Report",
        "slides": [
            {"title": "Summary", "bullets": ["AAPL up", "MSFT stable"], "image_paths": []},
            {"title": "Risks", "bullets": ["Volatility", "Macro"], "image_paths": []},
        ],
    }

    res = gen(
        path="out/weekly_report.pptx",
        deck=deck,
        workspace_id=TEST_WORKSPACE_ID,
        agent_id=TEST_AGENT_ID,
        session_id=TEST_SESSION_ID,
    )

    assert res.get("success") is True, res
    out_abs = res.get("output_path")
    assert out_abs, res
    prs = Presentation(str(out_abs))
    assert res.get("metadata", {}).get("slides") == 1 + len(deck["slides"])

    texts = []
    for s in prs.slides:
        for shape in s.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    joined ="\n".join(texts)
    assert "Weekly Report" in joined
    assert "Summary" in joined
    assert "Risks" in joined
    # ensure file exists in sandbox
    out_abs = (
        Path(powerpoint_tools["tmp_path"])
        / TEST_WORKSPACE_ID
        / TEST_AGENT_ID
        / TEST_SESSION_ID
        / "out"
        / "weekly_report.pptx"
    )
    assert out_abs.exists()
    assert out_abs.stat().st_size > 0


def test_powerpoint_generate_rejects_bad_extension(powerpoint_tools):
    gen = powerpoint_tools["powerpoint_generate"]
    res = gen(
        path="out/not_a_ppt.txt",
        deck={"title": "X", "slides": [{"title": "A", "bullets": []}]},
        workspace_id=TEST_WORKSPACE_ID,
        agent_id=TEST_AGENT_ID,
        session_id=TEST_SESSION_ID,
    )
    assert res.get("success") is False
    assert res.get("error") is not None

